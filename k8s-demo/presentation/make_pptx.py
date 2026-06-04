"""
Génère MotoShop_Soutenance.pptx
Lancer : python presentation/make_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy, os

# ── Palette ────────────────────────────────────────────────────
BG     = RGBColor(0x0f,0x17,0x2a)
CARD   = RGBColor(0x1e,0x29,0x3b)
BORDER = RGBColor(0x33,0x41,0x55)
BLUE   = RGBColor(0x3b,0x82,0xf6)
PURPLE = RGBColor(0x8b,0x5c,0xf6)
CYAN   = RGBColor(0x06,0xb6,0xd4)
GREEN  = RGBColor(0x10,0xb9,0x81)
ORANGE = RGBColor(0xf5,0x9e,0x0b)
RED    = RGBColor(0xef,0x44,0x44)
TEXT   = RGBColor(0xf1,0xf5,0xf9)
MUTED  = RGBColor(0x94,0xa3,0xb8)
WHITE  = RGBColor(0xff,0xff,0xff)
BLACK  = RGBColor(0x00,0x00,0x00)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]   # complètement vide

# ── Helpers ────────────────────────────────────────────────────

def add_slide():
    s = prs.slides.add_slide(blank)
    set_bg(s, BG)
    return s

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, fill=CARD, line=BORDER, line_w=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = line_w
    return shape

def accent_box(slide, x, y, w, h, color=BLUE):
    """Card avec barre colorée à gauche."""
    bg = box(slide, x, y, w, h, fill=CARD, line=BORDER)
    bar = box(slide, x, y, Inches(0.06), h, fill=color, line=color, line_w=Pt(0))
    return bg

def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf

def title_bar(slide, title, icon=""):
    """Barre de titre bleue standard."""
    box(slide, 0, 0, W, Inches(0.85), fill=CARD, line=BORDER)
    # Trait bleu en bas du header
    box(slide, 0, Inches(0.82), W, Inches(0.04), fill=BLUE, line=BLUE, line_w=Pt(0))
    txt(slide, f"{icon}  {title}" if icon else title,
        Inches(0.4), Inches(0.12), W - Inches(0.8), Inches(0.65),
        size=26, bold=True, color=TEXT)

def badge(slide, label, x, y, color=BLUE):
    w = Inches(1.5)
    h = Inches(0.32)
    b = box(slide, x, y, w, h, fill=RGBColor(
        int(color[0]*0.2), int(color[1]*0.2), int(color[2]*0.2)
    ), line=color, line_w=Pt(1.5))
    txt(slide, label, x, y + Inches(0.04), w, h - Inches(0.04),
        size=10, bold=True, color=color, align=PP_ALIGN.CENTER)
    return w + Inches(0.12)

def bullet_list(slide, items, x, y, w, size=14, color=TEXT, bullet_color=BLUE):
    cy = y
    for item in items:
        txt(slide, "▸", x, cy, Inches(0.25), Inches(0.32),
            size=size, bold=True, color=bullet_color)
        txt(slide, item, x + Inches(0.28), cy, w - Inches(0.28), Inches(0.38),
            size=size, color=color)
        cy += Inches(0.38)
    return cy

def stat_box(slide, number, label, x, y, w=Inches(2.8), color=BLUE):
    h = Inches(1.1)
    box(slide, x, y, w, h, fill=CARD, line=BORDER)
    txt(slide, number, x, y + Inches(0.05), w, Inches(0.62),
        size=30, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(slide, label.upper(), x, y + Inches(0.64), w, Inches(0.36),
        size=9, color=MUTED, align=PP_ALIGN.CENTER)

def section_title(slide, t, x, y, w=Inches(6), color=CYAN):
    txt(slide, t, x, y, w, Inches(0.32), size=14, bold=True, color=color)
    box(slide, x, y + Inches(0.32), Inches(1.2), Inches(0.03),
        fill=color, line=color, line_w=Pt(0))

def page_num(slide, n):
    txt(slide, str(n), W - Inches(0.6), H - Inches(0.35), Inches(0.4), Inches(0.28),
        size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITRE
# ══════════════════════════════════════════════════════════════
s = add_slide()

# Dégradé simulé avec deux rectangles
box(s, 0, 0, W, H * 0.55, fill=RGBColor(0x0d,0x1b,0x35), line=RGBColor(0x0d,0x1b,0x35), line_w=Pt(0))
box(s, 0, H * 0.55, W, H * 0.45, fill=BG, line=BG, line_w=Pt(0))

# Trait décoratif
box(s, Inches(4), Inches(0.18), Inches(5.5), Inches(0.05),
    fill=BLUE, line=BLUE, line_w=Pt(0))

# Icône moto
txt(s, "🏍️", Inches(6.1), Inches(0.9), Inches(1.2), Inches(1.2),
    size=54, align=PP_ALIGN.CENTER)

# Titre principal
txt(s, "MotoShop", Inches(2), Inches(2.0), Inches(9.3), Inches(1.4),
    size=72, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Sous-titre
txt(s, "Plateforme e-commerce déployée sur Kubernetes",
    Inches(2), Inches(3.3), Inches(9.3), Inches(0.55),
    size=22, color=MUTED, align=PP_ALIGN.CENTER)

# Badges technos
bx = Inches(2.5)
by = Inches(4.05)
for lbl, col in [("🐳  Docker",BLUE),("☸  Kubernetes",CYAN),
                  ("⚛  React 18",PURPLE),("🐍  Flask",GREEN),("🍃  MongoDB",ORANGE)]:
    bx += badge(s, lbl, bx, by, color=col)

# École + Date
txt(s, "Projet DevOps  —  BTS / BUT / Licence Informatique",
    Inches(2), Inches(4.75), Inches(9.3), Inches(0.4),
    size=14, color=CYAN, align=PP_ALIGN.CENTER, bold=True)
txt(s, "Alexandra Assaga  ·  Juin 2026",
    Inches(2), Inches(5.15), Inches(9.3), Inches(0.4),
    size=13, color=MUTED, align=PP_ALIGN.CENTER)

# Trait décoratif bas
box(s, Inches(4), H - Inches(0.22), Inches(5.5), Inches(0.05),
    fill=PURPLE, line=PURPLE, line_w=Pt(0))
page_num(s, 1)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — ÉQUIPE
# ══════════════════════════════════════════════════════════════
s = add_slide()
title_bar(s, "Présentation de l'équipe", "👥")

# Carte membre
accent_box(s, Inches(0.35), Inches(1.05), Inches(5.8), Inches(5.3), BLUE)
txt(s, "👩‍💻", Inches(0.55), Inches(1.2), Inches(0.9), Inches(0.9), size=36)
txt(s, "Alexandra Assaga",
    Inches(1.55), Inches(1.25), Inches(4.3), Inches(0.52),
    size=22, bold=True, color=TEXT)
txt(s, "Lead DevOps & Full-Stack",
    Inches(1.55), Inches(1.75), Inches(4.3), Inches(0.38),
    size=13, color=CYAN)

section_title(s, "Responsabilités", Inches(0.6), Inches(2.35))
items = [
    "Architecture Kubernetes & Docker",
    "Développement Backend — Flask / MongoDB",
    "Développement Frontend — React / Vite",
    "Pipeline CI/CD & scripts d'automatisation",
    "Documentation et déploiement GitHub",
]
bullet_list(s, items, Inches(0.6), Inches(2.75), Inches(5.2))

# Carte répartition tâches
accent_box(s, Inches(6.55), Inches(1.05), Inches(6.4), Inches(3.1), PURPLE)
section_title(s, "Répartition des tâches", Inches(6.8), Inches(1.2), color=PURPLE)

rows = [("Backend","API Flask + MongoDB + Routes"),
        ("Frontend","React SPA + Nginx + Pages"),
        ("Docker","Dockerfiles + Images Hub"),
        ("Kubernetes","9 Manifests YAML"),
        ("CI/CD","Hook git + deploy.ps1")]
ry = Inches(1.65)
for dom, task in rows:
    box(s, Inches(6.7), ry, Inches(1.7), Inches(0.38),
        fill=RGBColor(0x1a,0x23,0x3a), line=PURPLE, line_w=Pt(1))
    txt(s, dom, Inches(6.72), ry+Inches(0.05), Inches(1.66), Inches(0.32),
        size=11, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    txt(s, task, Inches(8.55), ry+Inches(0.05), Inches(4.2), Inches(0.32),
        size=11, color=TEXT)
    ry += Inches(0.46)

# Dépôt GitHub
accent_box(s, Inches(6.55), Inches(4.35), Inches(6.4), Inches(1.55), CYAN)
txt(s, "🔗  Dépôt GitHub",
    Inches(6.8), Inches(4.5), Inches(5.8), Inches(0.4),
    size=14, bold=True, color=CYAN)
txt(s, "github.com/aassaga21/MotoShop",
    Inches(6.8), Inches(4.95), Inches(5.8), Inches(0.4),
    size=13, color=TEXT)
txt(s, "hub.docker.com/u/alexandraassaga",
    Inches(6.8), Inches(5.33), Inches(5.8), Inches(0.4),
    size=13, color=MUTED)
page_num(s, 2)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — CONTEXTE & OBJECTIFS
# ══════════════════════════════════════════════════════════════
s = add_slide()
title_bar(s, "Contexte & Objectifs", "🎯")

cols = [
    (ORANGE, "🎯  Problématique", [
        "Déployer une app e-commerce complète",
        "Garantir la persistance des données",
        "Reproductibilité du déploiement",
        "Automatiser les mises à jour",
    ]),
    (BLUE, "🛒  Objectifs fonctionnels", [
        "Catalogue de 12 motos avec filtres",
        "Fiche produit détaillée",
        "Panier & formulaire de commande",
        "Page Contact (messages MongoDB)",
        "Dashboard Admin (commandes/clients)",
    ]),
    (CYAN, "⚙️  Objectifs techniques", [
        "Conteneurisation avec Docker",
        "Orchestration avec Kubernetes",
        "Persistance via PVC / PV",
        "Gestion des secrets K8s",
        "Pipeline CI/CD automatisé",
    ]),
]
cx = Inches(0.35)
for col, title, items in cols:
    w = Inches(4.15)
    accent_box(s, cx, Inches(1.05), w, Inches(4.5), col)
    txt(s, title, cx + Inches(0.18), Inches(1.15), w - Inches(0.3), Inches(0.45),
        size=14, bold=True, color=col)
    bullet_list(s, items, cx + Inches(0.18), Inches(1.65), w - Inches(0.3), size=13)
    cx += w + Inches(0.27)

# Stats
stats = [("12","Motos catalogue",BLUE),("3","Pods Kubernetes",PURPLE),
         ("5","Routes API",CYAN),("4","Pages React",GREEN)]
sx = Inches(0.35)
for num, lbl, col in stats:
    stat_box(s, num, lbl, sx, Inches(5.72), w=Inches(2.95), color=col)
    sx += Inches(3.1)
page_num(s, 3)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE GLOBALE
# ══════════════════════════════════════════════════════════════
s = add_slide()
title_bar(s, "Architecture Globale", "🌐")

def arch_node(slide, label, sublabel, x, y, w, h, color=BLUE):
    box(slide, x, y, w, h,
        fill=RGBColor(int(color[0]*0.15), int(color[1]*0.15), int(color[2]*0.15)),
        line=color, line_w=Pt(2))
    txt(slide, label, x, y + Inches(0.08), w, Inches(0.4),
        size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    if sublabel:
        txt(slide, sublabel, x, y + Inches(0.45), w, Inches(0.35),
            size=10, color=MUTED, align=PP_ALIGN.CENTER)

def arrow_down(slide, x, y, label=""):
    txt(slide, "▼", x, y, Inches(0.4), Inches(0.35),
        size=18, color=BLUE, align=PP_ALIGN.CENTER)
    if label:
        txt(slide, label, x + Inches(0.35), y + Inches(0.04), Inches(3), Inches(0.28),
            size=10, color=MUTED)

cx = Inches(4.5)

# Utilisateur
arch_node(s, "👤  Navigateur utilisateur", "http://localhost:30080", cx, Inches(1.0), Inches(4.3), Inches(0.75), MUTED)
arrow_down(s, cx + Inches(1.85), Inches(1.8), "  NodePort 30080")

# Pod 1 — Frontend
box(s, Inches(0.5), Inches(2.25), Inches(12.3), Inches(1.1), fill=RGBColor(0x0a,0x1a,0x33), line=BLUE, line_w=Pt(2))
txt(s, "POD 1  —  Frontend  (Namespace: ecommerce)", Inches(0.65), Inches(2.3), Inches(6), Inches(0.3),
    size=10, color=BLUE, bold=True)
arch_node(s, "🌐  Nginx + React 18", "Build statique Vite  ·  Port 80  ·  Proxy /api/ → backend-service:5000",
          Inches(2.5), Inches(2.5), Inches(8.3), Inches(0.75), BLUE)

arrow_down(s, cx + Inches(1.85), Inches(3.4), "  ClusterIP  ·  backend-service:5000")

# Pod 2 — Backend
box(s, Inches(0.5), Inches(3.85), Inches(12.3), Inches(1.35), fill=RGBColor(0x05,0x1a,0x12), line=GREEN, line_w=Pt(2))
txt(s, "POD 2  —  Backend + Base de données", Inches(0.65), Inches(3.9), Inches(6), Inches(0.3),
    size=10, color=GREEN, bold=True)
arch_node(s, "🐍  Flask API", "Port 5000  ·  Python 3.12",
          Inches(1.5), Inches(4.1), Inches(4.2), Inches(0.9), PURPLE)
txt(s, "⟺  localhost", Inches(5.85), Inches(4.3), Inches(1.5), Inches(0.4),
    size=14, color=MUTED, align=PP_ALIGN.CENTER)
arch_node(s, "🍃  MongoDB 7.0", "Port 27017  ·  PVC /data/db",
          Inches(7.5), Inches(4.1), Inches(4.2), Inches(0.9), GREEN)

# Secrets / Config / PVC
items = [("🔐  Secret K8s", "mongouser / MongoPass2024", ORANGE),
         ("⚙  ConfigMap", "MONGO_HOST · DB_NAME", CYAN),
         ("💾  PV / PVC", "1Gi · RWO · /data/db", GREEN)]
ix = Inches(1.0)
for lbl, sub, col in items:
    box(s, ix, Inches(5.35), Inches(3.7), Inches(0.75),
        fill=RGBColor(0x12,0x15,0x1a), line=col, line_w=Pt(1.5))
    txt(s, lbl, ix + Inches(0.12), Inches(5.4), Inches(3.4), Inches(0.32),
        size=11, bold=True, color=col)
    txt(s, sub, ix + Inches(0.12), Inches(5.7), Inches(3.4), Inches(0.32),
        size=10, color=MUTED)
    ix += Inches(4.0)
page_num(s, 4)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — ARCHITECTURE DOCKER
# ══════════════════════════════════════════════════════════════
s = add_slide()
title_bar(s, "Architecture Docker", "🐳")

# Frontend image
accent_box(s, Inches(0.35), Inches(1.05), Inches(6.0), Inches(2.6), BLUE)
txt(s, "🌐  Image Frontend", Inches(0.6), Inches(1.15), Inches(5.5), Inches(0.42),
    size=16, bold=True, color=BLUE)
txt(s, "alexandraassaga/moto-frontend:latest", Inches(0.6), Inches(1.55), Inches(5.5), Inches(0.32),
    size=11, color=CYAN)
items_fe = ["Base : node:20-alpine  →  nginx:alpine",
            "Build : npm run build  (Vite — multi-stage)",
            "Contient : SPA React + images statiques",
            "Port exposé : 80"]
bullet_list(s, items_fe, Inches(0.6), Inches(1.92), Inches(5.5), size=12)

# Backend image
accent_box(s, Inches(0.35), Inches(3.75), Inches(6.0), Inches(2.4), GREEN)
txt(s, "🐍  Image Backend", Inches(0.6), Inches(3.85), Inches(5.5), Inches(0.42),
    size=16, bold=True, color=GREEN)
txt(s, "alexandraassaga/moto-backend:latest", Inches(0.6), Inches(4.25), Inches(5.5), Inches(0.32),
    size=11, color=CYAN)
items_be = ["Base : python:3.12-slim",
            "Contient : Flask, PyMongo, Flask-CORS",
            "Port exposé : 5000",
            "Démarrage : python app.py"]
bullet_list(s, items_be, Inches(0.6), Inches(4.62), Inches(5.5), size=12)

# Volumes
accent_box(s, Inches(6.75), Inches(1.05), Inches(6.2), Inches(2.1), ORANGE)
txt(s, "💾  Volumes & Persistance", Inches(7.0), Inches(1.15), Inches(5.7), Inches(0.42),
    size=16, bold=True, color=ORANGE)
code_lines = [
    ("volumes:", CYAN),
    ("  - name: mongo-data", TEXT),
    ("    persistentVolumeClaim:", TEXT),
    ("      claimName: mongo-pvc", GREEN),
    ("# Montage : /data/db", MUTED),
    ("# Capacité : 1Gi  · ReadWriteOnce", MUTED),
]
cy = Inches(1.65)
for line, col in code_lines:
    txt(s, line, Inches(7.1), cy, Inches(5.5), Inches(0.32),
        size=11, color=col)
    cy += Inches(0.3)

# Avantages
accent_box(s, Inches(6.75), Inches(3.25), Inches(6.2), Inches(2.9), PURPLE)
txt(s, "✅  Avantages de la conteneurisation",
    Inches(7.0), Inches(3.35), Inches(5.7), Inches(0.42),
    size=16, bold=True, color=PURPLE)
items_adv = ["Isolation complète des dépendances",
             "Reproductibilité : même comportement Dev/Prod",
             "Images versionées sur Docker Hub",
             "Déploiement déclaratif via Kubernetes",
             "Mise à l'échelle simplifiée"]
bullet_list(s, items_adv, Inches(7.0), Inches(3.82), Inches(5.7), size=12)
page_num(s, 5)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — CHOIX TECHNIQUES
# ══════════════════════════════════════════════════════════════
s = add_slide()
title_bar(s, "Choix Techniques Justifiés", "⚙️")

techs = [
    ("🐳", "Docker", "Isolation · Portabilité\nImages reproductibles", BLUE),
    ("☸️", "Kubernetes", "Orchestration · Self-healing\nGestion des ressources", CYAN),
    ("🍃", "MongoDB 7.0", "NoSQL flexible · Schéma\ndynamique pour produits", GREEN),
    ("⚛️", "React 18 + Vite", "SPA moderne · Build rapide\nComposants réutilisables", PURPLE),
    ("🐍", "Python Flask", "API REST légère\nIntégration PyMongo native", GREEN),
    ("🔐", "Secrets K8s", "Credentials base64\nSéparation config / code", ORANGE),
]
gx, gy = Inches(0.35), Inches(1.05)
tw, th = Inches(4.1), Inches(2.65)
for i, (icon, name, desc, color) in enumerate(techs):
    col = i % 3
    row = i // 3
    x = gx + col * (tw + Inches(0.3))
    y = gy + row * (th + Inches(0.2))
    accent_box(s, x, y, tw, th, color)
    txt(s, icon, x + Inches(0.25), y + Inches(0.18), Inches(0.65), Inches(0.65), size=30)
    txt(s, name, x + Inches(0.18), y + Inches(0.85), tw - Inches(0.3), Inches(0.45),
        size=16, bold=True, color=color)
    txt(s, desc, x + Inches(0.18), y + Inches(1.32), tw - Inches(0.3), Inches(1.0),
        size=12, color=MUTED)
page_num(s, 6)


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — DÉMONSTRATION DOCKER
# ══════════════════════════════════════════════════════════════
s = add_slide()
title_bar(s, "Démonstration Docker", "🐳")

# Table conteneurs
accent_box(s, Inches(0.35), Inches(1.05), Inches(5.8), Inches(2.8), BLUE)
txt(s, "📦  Conteneurs en production",
    Inches(0.6), Inches(1.15), Inches(5.3), Inches(0.42), size=14, bold=True, color=BLUE)
headers = ["Conteneur", "Image", "Port", "Rôle"]
cols_w = [Inches(1.2), Inches(2.0), Inches(0.7), Inches(1.5)]
rows_data = [
    ("flask",    "moto-backend:latest",  "5000", "API REST"),
    ("mongodb",  "mongo:7.0",            "27017","Base de données"),
    ("frontend", "moto-frontend:latest", "80",   "SPA + Nginx"),
]
hx = Inches(0.6)
hy = Inches(1.65)
for i, h in enumerate(headers):
    box(s, hx, hy, cols_w[i], Inches(0.32), fill=RGBColor(0x0d,0x18,0x30), line=BORDER)
    txt(s, h, hx + Inches(0.05), hy + Inches(0.04), cols_w[i], Inches(0.26),
        size=10, bold=True, color=CYAN)
    hx += cols_w[i]
ry = Inches(1.97)
for row in rows_data:
    rx = Inches(0.6)
    for i, cell in enumerate(row):
        txt(s, cell, rx + Inches(0.05), ry + Inches(0.06), cols_w[i], Inches(0.3),
            size=11, color=TEXT)
        rx += cols_w[i]
    box(s, Inches(0.6), ry + Inches(0.35), Inches(5.45), Inches(0.02),
        fill=BORDER, line=BORDER, line_w=Pt(0))
    ry += Inches(0.42)

# Communication
accent_box(s, Inches(0.35), Inches(4.0), Inches(5.8), Inches(2.3), GREEN)
txt(s, "🔗  Communication inter-conteneurs",
    Inches(0.6), Inches(4.1), Inches(5.3), Inches(0.42), size=14, bold=True, color=GREEN)
for label, sublabel, col in [
    ("Nginx :80", "Frontend", BLUE),
    ("→  /api/", "", MUTED),
    ("Flask :5000", "Backend", PURPLE),
    ("→", "", MUTED),
    ("MongoDB :27017", "Sidecar", GREEN),
]:
    pass
# Schéma flux horizontal
fx = Inches(0.7)
fy = Inches(4.65)
for lbl, col in [("Nginx\n:80", BLUE),("→  /api/",MUTED),("Flask\n:5000",PURPLE),("→\nlocalhost",MUTED),("MongoDB\n:27017",GREEN)]:
    if "→" in lbl:
        txt(s, lbl.replace("\n"," "), fx, fy, Inches(0.9), Inches(0.7), size=10, color=col, align=PP_ALIGN.CENTER)
    else:
        box(s, fx, fy, Inches(1.2), Inches(0.7),
            fill=RGBColor(int(col[0]*0.1),int(col[1]*0.1),int(col[2]*0.1)),
            line=col, line_w=Pt(1.5))
        txt(s, lbl, fx, fy + Inches(0.1), Inches(1.2), Inches(0.55), size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    fx += Inches(1.1)
txt(s, "Flask et MongoDB partagent le même pod → communication via localhost (pas de réseau externe)",
    Inches(0.6), Inches(5.45), Inches(5.3), Inches(0.55), size=10, color=MUTED)

# Commandes Docker
accent_box(s, Inches(6.55), Inches(1.05), Inches(6.4), Inches(5.25), PURPLE)
txt(s, "💻  Commandes Docker utilisées",
    Inches(6.8), Inches(1.15), Inches(5.9), Inches(0.42), size=14, bold=True, color=PURPLE)
cmd_lines = [
    ("# Build image backend", MUTED),
    ("docker build -t alexandraassaga/moto-backend:latest ./backend", CYAN),
    ("", TEXT),
    ("# Push sur Docker Hub", MUTED),
    ("docker push alexandraassaga/moto-backend:latest", CYAN),
    ("", TEXT),
    ("# Build multi-stage frontend", MUTED),
    ("docker build -t alexandraassaga/moto-frontend:latest ./frontend", CYAN),
    ("", TEXT),
    ("# Vérifier les pods K8s", MUTED),
    ("kubectl get pods -n ecommerce", GREEN),
    ("", TEXT),
    ("# Logs du container Flask", MUTED),
    ("kubectl logs deployment/backend -c flask -n ecommerce", GREEN),
]
cy = Inches(1.65)
for line, col in cmd_lines:
    if line:
        txt(s, line, Inches(6.9), cy, Inches(6.1), Inches(0.3), size=10, color=col)
    cy += Inches(0.3)
page_num(s, 7)


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — CLUSTER KUBERNETES
# ══════════════════════════════════════════════════════════════
s = add_slide()
title_bar(s, "Architecture du Cluster Kubernetes", "☸")

# Environnement
accent_box(s, Inches(0.35), Inches(1.05), Inches(5.8), Inches(2.7), CYAN)
txt(s, "🖥️  Environnement", Inches(0.6), Inches(1.15), Inches(5.3), Inches(0.42),
    size=16, bold=True, color=CYAN)
items_env = ["Cluster : Docker Desktop (single-node)",
             "Version Kubernetes : 1.34",
             "Namespace : ecommerce",
             "Nœud : docker-desktop  (IP 192.168.65.3)",
             "Accès : http://localhost:30080"]
bullet_list(s, items_env, Inches(0.6), Inches(1.62), Inches(5.3), size=12)

# Ressources
accent_box(s, Inches(0.35), Inches(3.85), Inches(5.8), Inches(2.4), PURPLE)
txt(s, "📊  Ressources allouées", Inches(0.6), Inches(3.95), Inches(5.3), Inches(0.42),
    size=16, bold=True, color=PURPLE)
headers2 = ["Conteneur", "RAM (req→limit)", "CPU (req→limit)"]
rows2 = [("flask","128Mi → 256Mi","100m → 300m"),
         ("mongodb","256Mi → 512Mi","200m → 500m"),
         ("frontend","64Mi → 128Mi","50m → 100m")]
hx = Inches(0.55)
hy = Inches(4.42)
for h, w in zip(headers2, [Inches(1.4), Inches(2.3), Inches(1.8)]):
    box(s, hx, hy, w, Inches(0.3), fill=RGBColor(0x1a,0x0a,0x30), line=BORDER)
    txt(s, h, hx+Inches(0.05), hy+Inches(0.03), w, Inches(0.26), size=9, bold=True, color=PURPLE)
    hx += w
ry = Inches(4.72)
for row in rows2:
    rx = Inches(0.55)
    for cell, w in zip(row, [Inches(1.4), Inches(2.3), Inches(1.8)]):
        txt(s, cell, rx+Inches(0.05), ry+Inches(0.04), w, Inches(0.3), size=10, color=TEXT)
        rx += w
    ry += Inches(0.38)

# Objets K8s
accent_box(s, Inches(6.55), Inches(1.05), Inches(6.4), Inches(5.25), BLUE)
txt(s, "🏗️  Objets Kubernetes déployés",
    Inches(6.8), Inches(1.15), Inches(5.9), Inches(0.42), size=16, bold=True, color=BLUE)

objs = [("Deployment","backend · frontend", BLUE),
        ("Service","ClusterIP + NodePort", CYAN),
        ("ConfigMap","app-config", PURPLE),
        ("Secret","mongo-secret (base64)", ORANGE),
        ("PersistentVolume","mongo-pv · storageClass: manual", GREEN),
        ("PersistentVolumeClaim","mongo-pvc · 1Gi · ReadWriteOnce", GREEN)]
oy = Inches(1.65)
for kind, detail, col in objs:
    box(s, Inches(6.7), oy, Inches(1.9), Inches(0.52),
        fill=RGBColor(int(col[0]*0.1), int(col[1]*0.1), int(col[2]*0.1)), line=col, line_w=Pt(1.5))
    txt(s, kind, Inches(6.72), oy+Inches(0.06), Inches(1.86), Inches(0.44),
        size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, detail, Inches(8.75), oy+Inches(0.1), Inches(4.0), Inches(0.38), size=11, color=TEXT)
    oy += Inches(0.65)
page_num(s, 8)


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — MANIFESTS YAML
# ══════════════════════════════════════════════════════════════
s = add_slide()
title_bar(s, "Déploiement Kubernetes — Manifests YAML", "📄")

# backend-deployment.yaml
accent_box(s, Inches(0.35), Inches(1.05), Inches(6.0), Inches(5.25), BLUE)
txt(s, "backend-deployment.yaml  (extrait clé)",
    Inches(0.6), Inches(1.15), Inches(5.5), Inches(0.38), size=13, bold=True, color=BLUE)
yaml_lines_be = [
    ("apiVersion: apps/v1", CYAN),
    ("kind: Deployment", TEXT),
    ("spec:", TEXT),
    ("  replicas: 1", TEXT),
    ("  strategy:", TEXT),
    ("    type: Recreate    # crucial !", GREEN),
    ("  containers:", TEXT),
    ("  - name: flask", PURPLE),
    ("    image: alexandraassaga/moto-backend", TEXT),
    ("    resources:", TEXT),
    ("      limits:", TEXT),
    ("        memory: \"256Mi\"", ORANGE),
    ("        cpu: \"300m\"", ORANGE),
    ("  - name: mongodb", GREEN),
    ("    image: mongo:7.0", TEXT),
    ("    volumeMounts:", TEXT),
    ("    - mountPath: /data/db", CYAN),
]
cy = Inches(1.6)
for line, col in yaml_lines_be:
    txt(s, line, Inches(0.55), cy, Inches(5.5), Inches(0.28), size=10, color=col)
    cy += Inches(0.28)

# Secrets + PVC
for title, lines, col, gy in [
    ("secret.yaml", [
        ("kind: Secret", TEXT),
        ("type: Opaque", TEXT),
        ("data:", TEXT),
        ("  MONGO_USERNAME: bW9uZ291c2Vy", ORANGE),
        ("  MONGO_PASSWORD: TW9uZ29QYXNz...", ORANGE),
    ], ORANGE, Inches(1.05)),
    ("mongo-pvc.yaml", [
        ("kind: PersistentVolumeClaim", TEXT),
        ("spec:", TEXT),
        ("  accessModes:", TEXT),
        ("  - ReadWriteOnce", GREEN),
        ("  resources:", TEXT),
        ("    requests:", TEXT),
        ("      storage: 1Gi", GREEN),
    ], GREEN, Inches(3.05)),
]:
    accent_box(s, Inches(6.55), gy, Inches(6.4), Inches(1.85), col)
    txt(s, title, Inches(6.8), gy + Inches(0.1), Inches(5.9), Inches(0.35),
        size=12, bold=True, color=col)
    cy = gy + Inches(0.5)
    for line, lc in lines:
        txt(s, line, Inches(6.8), cy, Inches(5.9), Inches(0.27), size=10, color=lc)
        cy += Inches(0.27)

# Note importante
box(s, Inches(6.55), Inches(5.05), Inches(6.4), Inches(1.25),
    fill=RGBColor(0x1a,0x12,0x00), line=ORANGE, line_w=Pt(2))
txt(s, "⚠️  Pourquoi strategy: Recreate ?",
    Inches(6.8), Inches(5.12), Inches(5.9), Inches(0.38), size=12, bold=True, color=ORANGE)
txt(s, "Flask et MongoDB partagent un PVC ReadWriteOnce.\nAvec RollingUpdate, 2 pods montent le même volume → conflit mongod.lock.",
    Inches(6.8), Inches(5.5), Inches(5.9), Inches(0.7), size=11, color=TEXT)
page_num(s, 9)


# ══════════════════════════════════════════════════════════════
# SLIDE 10 — PODS & SERVICES
# ══════════════════════════════════════════════════════════════
s = add_slide()
title_bar(s, "Pods & Services Kubernetes", "🔌")

# Pod Frontend
box(s, Inches(0.35), Inches(1.05), Inches(5.8), Inches(1.5),
    fill=RGBColor(0x05,0x10,0x25), line=BLUE, line_w=Pt(2))
txt(s, "POD — frontend", Inches(0.5), Inches(1.1), Inches(3), Inches(0.3),
    size=10, bold=True, color=BLUE)
box(s, Inches(0.8), Inches(1.42), Inches(5.0), Inches(0.9),
    fill=RGBColor(0x08,0x18,0x35), line=BLUE, line_w=Pt(1.5))
txt(s, "🌐  nginx:alpine", Inches(0.8), Inches(1.52), Inches(5.0), Inches(0.38),
    size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
txt(s, "Port 80  ·  SPA React statique", Inches(0.8), Inches(1.9), Inches(5.0), Inches(0.3),
    size=10, color=MUTED, align=PP_ALIGN.CENTER)

# Pod Backend
box(s, Inches(0.35), Inches(2.75), Inches(5.8), Inches(2.1),
    fill=RGBColor(0x08,0x18,0x10), line=GREEN, line_w=Pt(2))
txt(s, "POD — backend  (2 conteneurs — même pod)", Inches(0.5), Inches(2.8), Inches(5.5), Inches(0.3),
    size=10, bold=True, color=GREEN)
for cx2, name, port, col in [(Inches(0.55),  "🐍  Flask\nAPI REST",    "Port 5000", PURPLE),
                               (Inches(3.05),  "🍃  MongoDB\nBase de données", "Port 27017", GREEN)]:
    box(s, cx2, Inches(3.15), Inches(2.3), Inches(1.5),
        fill=RGBColor(int(col[0]*0.08),int(col[1]*0.08),int(col[2]*0.08)), line=col, line_w=Pt(1.5))
    txt(s, name, cx2, Inches(3.28), Inches(2.3), Inches(0.65), size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, port, cx2, Inches(3.95), Inches(2.3), Inches(0.3), size=10, color=MUTED, align=PP_ALIGN.CENTER)
txt(s, "⟺  localhost", Inches(2.95), Inches(3.6), Inches(0.9), Inches(0.38),
    size=12, color=MUTED, align=PP_ALIGN.CENTER)

# Services
accent_box(s, Inches(0.35), Inches(4.95), Inches(5.8), Inches(1.4), CYAN)
txt(s, "🔗  Services exposés",
    Inches(0.6), Inches(5.05), Inches(5.3), Inches(0.38), size=14, bold=True, color=CYAN)
for sx2, svc, stype, port, col in [
    (Inches(0.6),  "frontend-service", "NodePort", "30080 → 80",   ORANGE),
    (Inches(3.25), "backend-service",  "ClusterIP","5000 → 5000", CYAN),
]:
    box(s, sx2, Inches(5.5), Inches(2.5), Inches(0.65),
        fill=RGBColor(0x10,0x14,0x22), line=col, line_w=Pt(1.5))
    txt(s, svc, sx2 + Inches(0.08), Inches(5.55), Inches(2.3), Inches(0.28),
        size=11, bold=True, color=col)
    txt(s, f"{stype}  ·  {port}", sx2 + Inches(0.08), Inches(5.82), Inches(2.3), Inches(0.28),
        size=10, color=TEXT)

# Découverte services
accent_box(s, Inches(6.55), Inches(1.05), Inches(6.4), Inches(5.3), PURPLE)
txt(s, "📡  Découverte de services & Communication",
    Inches(6.8), Inches(1.15), Inches(5.9), Inches(0.42), size=14, bold=True, color=PURPLE)
items_svc = [
    "Nginx proxifie /api/ vers backend-service:5000",
    "Flask accède à MongoDB via localhost:27017",
    "(même pod = réseau partagé, pas de latence réseau)",
    "DNS K8s interne : backend-service.ecommerce.svc.cluster.local",
    "Accès externe navigateur : http://localhost:30080",
    "readinessProbe Flask : GET /api/health (delay 30s)",
]
bullet_list(s, items_svc, Inches(6.8), Inches(1.65), Inches(5.9), size=12)

# Schéma flux
box(s, Inches(6.7), Inches(4.35), Inches(6.1), Inches(1.85),
    fill=RGBColor(0x08,0x10,0x22), line=BORDER)
fx2 = Inches(6.9)
for lbl, col in [("Navigateur",MUTED),("→",MUTED),("NodePort\n30080",ORANGE),("→",MUTED),
                  ("Nginx\n:80",BLUE),("→ /api/",MUTED),("Flask\n:5000",PURPLE),("→",MUTED),("MongoDB\n:27017",GREEN)]:
    if "→" in lbl:
        txt(s, lbl.replace("\n"," "), fx2, Inches(4.7), Inches(0.65), Inches(0.5), size=9, color=col, align=PP_ALIGN.CENTER)
        fx2 += Inches(0.55)
    else:
        box(s, fx2, Inches(4.55), Inches(1.05), Inches(0.8),
            fill=RGBColor(int(col[0]*0.08),int(col[1]*0.08),int(col[2]*0.08)), line=col, line_w=Pt(1))
        txt(s, lbl, fx2, Inches(4.62), Inches(1.05), Inches(0.65), size=9, bold=True, color=col, align=PP_ALIGN.CENTER)
        fx2 += Inches(1.1)
page_num(s, 10)


# ══════════════════════════════════════════════════════════════
# SLIDE 11 — POD MONITORING
# ══════════════════════════════════════════════════════════════
s = add_slide()
title_bar(s, "Pod Monitoring — Prometheus & Grafana", "📊")

# Encadré du pod
box(s, Inches(0.35), Inches(1.05), Inches(5.8), Inches(5.3),
    fill=RGBColor(0x08,0x10,0x1a), line=ORANGE, line_w=Pt(2))
txt(s, "POD — monitoring-pod  (Namespace: ecommerce)",
    Inches(0.5), Inches(1.1), Inches(5.5), Inches(0.3),
    size=10, bold=True, color=ORANGE)

# Container Prometheus
box(s, Inches(0.55), Inches(1.48), Inches(5.0), Inches(1.7),
    fill=RGBColor(0x18,0x10,0x00), line=ORANGE, line_w=Pt(1.5))
txt(s, "📡  Prometheus v2.52", Inches(0.55), Inches(1.56), Inches(5.0), Inches(0.38),
    size=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
bullet_list(s, ["Collecte & stocke les métriques","Scrape les pods du namespace ecommerce","Port 9090  ·  NodePort 30900"],
    Inches(0.78), Inches(2.0), Inches(4.6), size=11, bullet_color=ORANGE)

# Séparateur
txt(s, "⟺  localhost  (réseau partagé du pod)",
    Inches(0.9), Inches(3.24), Inches(4.5), Inches(0.35),
    size=11, color=MUTED, align=PP_ALIGN.CENTER)

# Container Grafana
box(s, Inches(0.55), Inches(3.62), Inches(5.0), Inches(1.7),
    fill=RGBColor(0x04,0x14,0x04), line=GREEN, line_w=Pt(1.5))
txt(s, "📈  Grafana 10.4", Inches(0.55), Inches(3.7), Inches(5.0), Inches(0.38),
    size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
bullet_list(s, ["Dashboards de visualisation","Datasource : http://localhost:9090","Port 3000  ·  NodePort 30300"],
    Inches(0.78), Inches(4.14), Inches(4.6), size=11, bullet_color=GREEN)

# Pourquoi 1 pod
accent_box(s, Inches(6.55), Inches(1.05), Inches(6.4), Inches(2.55), CYAN)
txt(s, "⚡  Pourquoi 1 seul pod ?",
    Inches(6.8), Inches(1.15), Inches(5.9), Inches(0.38), size=14, bold=True, color=CYAN)
bullet_list(s, [
    "Grafana ↔ Prometheus via localhost (0 latence)",
    "Pas de DNS inter-service nécessaire",
    "ConfigMap monte prometheus.yml dans Prometheus",
    "emptyDir pour stockage temporaire des métriques",
], Inches(6.8), Inches(1.6), Inches(5.9), size=12)

# Commandes
accent_box(s, Inches(6.55), Inches(3.7), Inches(6.4), Inches(2.65), GREEN)
txt(s, "🌐  Déploiement & Accès",
    Inches(6.8), Inches(3.8), Inches(5.9), Inches(0.38), size=14, bold=True, color=GREEN)
for line, col in [
    ("# Appliquer les 3 manifests", MUTED),
    ("kubectl apply -f k8s/monitoring-configmap.yaml", CYAN),
    ("kubectl apply -f k8s/monitoring-pod.yaml", CYAN),
    ("kubectl apply -f k8s/monitoring-service.yaml", CYAN),
    ("", TEXT),
    ("# Accès (port-forward recommandé)", MUTED),
    ("kubectl port-forward pod/monitoring-pod 3000:3000 -n ecommerce", GREEN),
    ("# Grafana → http://localhost:3000  (admin / admin123)", MUTED),
    ("# Prometheus → http://localhost:9090", MUTED),
]:
    if line:
        txt(s, line, Inches(6.8), Inches(4.3), Inches(6.1), Inches(0.28), size=9, color=col)
    Inches(4.3)  # placeholder — loop variable needed

# Rebuild propre des lignes de commandes
cmd_lines_mon = [
    ("# Appliquer les 3 manifests", MUTED),
    ("kubectl apply -f k8s/monitoring-configmap.yaml", CYAN),
    ("kubectl apply -f k8s/monitoring-pod.yaml", CYAN),
    ("kubectl apply -f k8s/monitoring-service.yaml", CYAN),
    ("", TEXT),
    ("# Accès Grafana (admin / admin123)", MUTED),
    ("kubectl port-forward pod/monitoring-pod 3000:3000 \\", GREEN),
    ("  -n ecommerce   →  http://localhost:3000", GREEN),
    ("# Prometheus  →  http://localhost:9090", MUTED),
]
cy_mon = Inches(4.28)
for line, col in cmd_lines_mon:
    if line:
        txt(s, line, Inches(6.8), cy_mon, Inches(6.1), Inches(0.27), size=9, color=col)
    cy_mon += Inches(0.27)

page_num(s, 11)


# ══════════════════════════════════════════════════════════════
# SLIDE 12 — DÉMONSTRATION FONCTIONNELLE (ex-11)
# ══════════════════════════════════════════════════════════════
s = add_slide()
title_bar(s, "Démonstration Fonctionnelle", "▶")

steps = [
    ("1","Accès au catalogue","http://localhost:30080 — 12 motos avec filtres\nNaked · Sport · Trail · Cruiser", BLUE),
    ("2","Fiche produit","Clic sur une moto → modal :\nmoteur, puissance, poids, vitesse max, prix", CYAN),
    ("3","Panier & commande","Ajout au panier → sidebar → formulaire client\nPOST /api/orders → toast de confirmation", PURPLE),
    ("4","Page Contact","Formulaire (Nom/Email/Sujet/Message)\nPOST /api/contact → sauvegarde MongoDB", ORANGE),
    ("5","Dashboard Admin  ⚙","Stats CA · Gestion statuts commandes\nHistorique complet par client", GREEN),
]
sy = Inches(1.05)
for num, title, desc, col in steps:
    box(s, Inches(0.35), sy, Inches(0.48), Inches(0.75),
        fill=col, line=col, line_w=Pt(0))
    txt(s, num, Inches(0.35), sy + Inches(0.12), Inches(0.48), Inches(0.5),
        size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, Inches(0.85), sy, Inches(5.8), Inches(0.75),
        fill=RGBColor(int(col[0]*0.08),int(col[1]*0.08),int(col[2]*0.08)), line=col, line_w=Pt(1))
    txt(s, title, Inches(1.05), sy + Inches(0.05), Inches(5.5), Inches(0.34),
        size=13, bold=True, color=col)
    txt(s, desc, Inches(1.05), sy + Inches(0.38), Inches(5.5), Inches(0.55),
        size=10, color=MUTED)
    sy += Inches(0.95)

# Routes API
accent_box(s, Inches(6.95), Inches(1.05), Inches(6.0), Inches(3.3), CYAN)
txt(s, "📡  Routes API Flask",
    Inches(7.2), Inches(1.15), Inches(5.5), Inches(0.42), size=14, bold=True, color=CYAN)
routes = [
    ("GET",   "/api/products",           "Catalogue motos"),
    ("GET",   "/api/products/:id",       "Fiche produit"),
    ("POST",  "/api/orders",             "Créer commande"),
    ("GET",   "/api/orders",             "Admin : liste"),
    ("PATCH", "/api/orders/:id/status",  "Changer statut"),
    ("POST",  "/api/contact",            "Message contact"),
    ("GET",   "/api/health",             "Health check K8s"),
]
meth_colors = {"GET": GREEN, "POST": BLUE, "PATCH": ORANGE, "DELETE": RED}
ry = Inches(1.65)
for meth, route, desc in routes:
    mc = meth_colors.get(meth, CYAN)
    box(s, Inches(7.1), ry, Inches(0.72), Inches(0.34),
        fill=RGBColor(int(mc[0]*0.1),int(mc[1]*0.1),int(mc[2]*0.1)), line=mc, line_w=Pt(1))
    txt(s, meth, Inches(7.1), ry + Inches(0.04), Inches(0.72), Inches(0.3),
        size=9, bold=True, color=mc, align=PP_ALIGN.CENTER)
    txt(s, route, Inches(7.9), ry + Inches(0.04), Inches(2.4), Inches(0.3), size=10, color=TEXT)
    txt(s, desc, Inches(10.4), ry + Inches(0.04), Inches(2.4), Inches(0.3), size=10, color=MUTED)
    ry += Inches(0.4)

# Persistance
accent_box(s, Inches(6.95), Inches(4.45), Inches(6.0), Inches(1.85), GREEN)
txt(s, "✅  Test de persistance",
    Inches(7.2), Inches(4.55), Inches(5.5), Inches(0.38), size=14, bold=True, color=GREEN)
items_p = ["Suppression du pod backend → recréé automatiquement",
           "Données MongoDB conservées grâce au PVC (1Gi)",
           "Flask re-seed si collection vide (12 produits)"]
bullet_list(s, items_p, Inches(7.2), Inches(5.0), Inches(5.5), size=12)
page_num(s, 12)


# ══════════════════════════════════════════════════════════════
# SLIDE 13 — DIFFICULTÉS
# ══════════════════════════════════════════════════════════════
s = add_slide()
title_bar(s, "Difficultés Rencontrées", "⚠")

difficulties = [
    ("🔒","Conflit mongod.lock (CrashLoopBackOff)",
     "RollingUpdate crée un nouveau pod avant de supprimer l'ancien → 2 pods montent le même PVC → verrou bloqué",
     RED),
    ("💥","OOM Kill — exit code 137",
     "Lancer mongosh dans le conteneur MongoDB (512Mi) en plus de mongod dépasse la limite mémoire → SIGKILL",
     ORANGE),
    ("👤","User MongoDB inexistant (UserNotFound)",
     "MONGO_INITDB_* ne fonctionne que sur un volume VIDE. rm -rf /data/db/* ne supprime pas les fichiers cachés",
     ORANGE),
    ("🖼️","Images motos invisibles",
     "Noms de fichiers avec majuscules/espaces (Windows) — Linux est sensible à la casse → Nginx retourne 404",
     CYAN),
    ("🔗","Hook git installé au mauvais endroit",
     "Racine git à C:\\Users\\alexa\\ (parent de k8s-demo) → hook créé dans un .git\\ inexistant",
     PURPLE),
    ("⏱️","Race condition Flask / MongoDB",
     "Flask tente la connexion avant que MongoDB soit prêt → boucle d'erreurs d'authentification (30 tentatives)",
     BLUE),
]
for i, (icon, title, desc, col) in enumerate(difficulties):
    col_n = i % 2
    row_n = i // 2
    x = Inches(0.35) + col_n * Inches(6.55)
    y = Inches(1.05) + row_n * Inches(1.85)
    box(s, x, y, Inches(6.2), Inches(1.75),
        fill=RGBColor(int(col[0]*0.08),int(col[1]*0.08),int(col[2]*0.08)), line=col, line_w=Pt(2))
    txt(s, icon + "  " + title, x + Inches(0.15), y + Inches(0.1), Inches(5.9), Inches(0.42),
        size=13, bold=True, color=col)
    txt(s, desc, x + Inches(0.15), y + Inches(0.55), Inches(5.9), Inches(1.0),
        size=11, color=TEXT)
page_num(s, 13)


# ══════════════════════════════════════════════════════════════
# SLIDE 14 — SOLUTIONS
# ══════════════════════════════════════════════════════════════
s = add_slide()
title_bar(s, "Solutions Apportées", "✅")

solutions = [
    ("✅","strategy: Recreate",
     "Ajouté dans le Deployment backend → ancien pod supprimé AVANT le nouveau → PVC libéré",
     GREEN),
    ("✅","Bypass via container Flask",
     "python3 -c avec PyMongo déjà chargé au lieu de mongosh → 0 mémoire supplémentaire",
     GREEN),
    ("✅","Création manuelle du user MongoDB",
     "Pod temporaire mongod --noauth + getSiblingDB('admin').createUser() (--eval incompatible avec use admin)",
     CYAN),
    ("✅","Convention kebab-case minuscules",
     "Renommage des images (bmw-r1250gs.jpg) + rebuild Docker frontend pour inclure les nouveaux noms",
     CYAN),
    ("✅","git rev-parse --show-toplevel",
     "Détection automatique de la vraie racine git dans install-hook.ps1 → hook installé au bon endroit",
     PURPLE),
    ("✅","find -mindepth 1 -delete",
     "Supprime TOUS les fichiers y compris les cachés → MONGO_INITDB_* peut s'exécuter sur volume vide",
     ORANGE),
]
for i, (icon, title, desc, col) in enumerate(solutions):
    col_n = i % 2
    row_n = i // 2
    x = Inches(0.35) + col_n * Inches(6.55)
    y = Inches(1.05) + row_n * Inches(1.85)
    box(s, x, y, Inches(6.2), Inches(1.75),
        fill=RGBColor(int(col[0]*0.06),int(col[1]*0.06),int(col[2]*0.06)), line=col, line_w=Pt(2))
    txt(s, icon + "  " + title, x + Inches(0.15), y + Inches(0.1), Inches(5.9), Inches(0.42),
        size=13, bold=True, color=col)
    txt(s, desc, x + Inches(0.15), y + Inches(0.55), Inches(5.9), Inches(1.0),
        size=11, color=TEXT)
page_num(s, 14)


# ══════════════════════════════════════════════════════════════
# SLIDE 15 — BILAN
# ══════════════════════════════════════════════════════════════
s = add_slide()
title_bar(s, "Bilan du Projet", "📊")

for x, title, items, col in [
    (Inches(0.35), "🎯  Objectifs atteints", [
        "Application e-commerce fonctionnelle",
        "2/2 pods Running en production",
        "Persistance MongoDB validée (PVC)",
        "CI/CD via hook git automatisé",
        "Pages Contact & Admin opérationnelles",
        "Code versioné sur GitHub",
    ], GREEN),
    (Inches(4.7), "📚  Compétences acquises", [
        "Conteneurisation Docker avancée",
        "Orchestration Kubernetes (9 manifests)",
        "Gestion Secrets / ConfigMaps",
        "Débogage de pods en production",
        "Scripting PowerShell & bash",
        "Pipeline CI/CD local",
    ], BLUE),
    (Inches(9.05), "💡  Valeur ajoutée", [
        "Architecture microservices réelle",
        "Pattern sidecar Flask + MongoDB",
        "Workflow DevOps de bout en bout",
        "Documentation README exhaustive",
        "Présentation HTML + PowerPoint",
    ], PURPLE),
]:
    accent_box(s, x, Inches(1.05), Inches(4.1), Inches(4.6), col)
    txt(s, title, x + Inches(0.18), Inches(1.15), Inches(3.8), Inches(0.45),
        size=14, bold=True, color=col)
    bullet_list(s, items, x + Inches(0.18), Inches(1.68), Inches(3.85), size=12)

# Stats finales
stats14 = [("100%","Pods Running",GREEN),("12","Manifests YAML",BLUE),
           ("5","Routes API",CYAN),("Auto","CI/CD",PURPLE)]
sx = Inches(0.35)
for num, lbl, col in stats14:
    stat_box(s, num, lbl, sx, Inches(5.78), w=Inches(3.08), color=col)
    sx += Inches(3.23)
page_num(s, 15)


# ══════════════════════════════════════════════════════════════
# SLIDE 16 — PERSPECTIVES
# ══════════════════════════════════════════════════════════════
s = add_slide()
title_bar(s, "Perspectives d'Amélioration", "🚀")

perspectives = [
    (BLUE,   "🔄  CI/CD avancé",      ["Runner GitHub Actions cloud","Tests Pytest / Jest automatisés","Analyse SonarQube","Versionning sémantique des images"]),
    (ORANGE, "📊  Monitoring",         ["✓ Prometheus + Grafana (implémenté)","Alerting automatique sur les pods","Métriques MongoDB","Logs centralisés (ELK Stack)"]),
    (CYAN,   "📈  Scalabilité",        ["Horizontal Pod Autoscaler","MongoDB ReplicaSet","Cluster multi-nœuds","Ingress Nginx (Load Balancer)"]),
    (RED,    "🔐  Sécurité renforcée", ["HTTPS / TLS (cert-manager)","Authentification JWT","Network Policies K8s","HashiCorp Vault pour secrets"]),
    (PURPLE, "☁️  Déploiement Cloud",  ["GKE / EKS / AKS","Infrastructure as Code (Terraform)","CDN pour assets statiques","Backup MongoDB automatisé"]),
    (GREEN,  "🛠️  Fonctionnalités",    ["Authentification utilisateurs","Paiement Stripe","Notifications email","Application mobile React Native"]),
]
gx, gy = Inches(0.35), Inches(1.05)
tw2, th2 = Inches(4.1), Inches(2.65)
for i, (col, title, items) in enumerate(perspectives):
    c = i % 3
    r = i // 3
    x = gx + c * (tw2 + Inches(0.3))
    y = gy + r * (th2 + Inches(0.2))
    accent_box(s, x, y, tw2, th2, col)
    txt(s, title, x + Inches(0.18), y + Inches(0.1), tw2 - Inches(0.3), Inches(0.42),
        size=13, bold=True, color=col)
    bullet_list(s, items, x + Inches(0.18), y + Inches(0.6), tw2 - Inches(0.3), size=11)
page_num(s, 16)


# ══════════════════════════════════════════════════════════════
# SLIDE 17 — QUESTIONS / REMERCIEMENTS
# ══════════════════════════════════════════════════════════════
s = add_slide()

# Fond dégradé
box(s, 0, 0, W, H * 0.5,
    fill=RGBColor(0x0d,0x1b,0x35), line=RGBColor(0x0d,0x1b,0x35), line_w=Pt(0))
# Trait décoratif
box(s, Inches(3.5), Inches(0.2), Inches(6.3), Inches(0.06),
    fill=BLUE, line=BLUE, line_w=Pt(0))

txt(s, "🙏", Inches(5.8), Inches(0.55), Inches(1.7), Inches(1.3), size=60, align=PP_ALIGN.CENTER)
txt(s, "Merci pour votre attention",
    Inches(1.5), Inches(1.75), Inches(10.3), Inches(0.9),
    size=40, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
txt(s, "Projet DevOps  —  MotoShop E-Commerce Kubernetes",
    Inches(1.5), Inches(2.65), Inches(10.3), Inches(0.45),
    size=16, color=MUTED, align=PP_ALIGN.CENTER)

# Badges
bx = Inches(2.6); by = Inches(3.2)
for lbl, col in [("🐳  Docker",BLUE),("☸  Kubernetes",CYAN),("⚛  React 18",PURPLE),("🐍  Flask",GREEN),("🍃  MongoDB",ORANGE)]:
    bx += badge(s, lbl, bx, by, color=col)

# Ressources
accent_box(s, Inches(3.0), Inches(4.05), Inches(7.3), Inches(1.65), CYAN)
txt(s, "📎  Ressources du projet",
    Inches(3.3), Inches(4.15), Inches(6.8), Inches(0.38), size=14, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
txt(s, "🐙  github.com/aassaga21/MotoShop",
    Inches(3.3), Inches(4.6), Inches(6.8), Inches(0.32), size=13, color=TEXT, align=PP_ALIGN.CENTER)
txt(s, "🐳  hub.docker.com/u/alexandraassaga",
    Inches(3.3), Inches(4.93), Inches(6.8), Inches(0.32), size=13, color=MUTED, align=PP_ALIGN.CENTER)

# Box Questions
box(s, Inches(3.5), Inches(5.85), Inches(6.3), Inches(0.9),
    fill=RGBColor(0x08,0x15,0x2c), line=BLUE, line_w=Pt(3))
txt(s, "❓  Questions du jury",
    Inches(3.5), Inches(5.97), Inches(6.3), Inches(0.6),
    size=22, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

# Trait décoratif bas
box(s, Inches(3.5), H - Inches(0.22), Inches(6.3), Inches(0.06),
    fill=PURPLE, line=PURPLE, line_w=Pt(0))
page_num(s, 17)


# ── Sauvegarder ────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "MotoShop_Soutenance_v2.pptx")
prs.save(out)
print(f"\nFichier genere : {out}")
print(f"    {prs.slides.__len__()} diapositives - Format 16:9")
